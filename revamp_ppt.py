from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path(r"C:\Users\Saksh\Desktop\Network_Layer_IP_Addressing.pptx")
OUTPUT = Path(r"C:\Users\Saksh\Desktop\CampusSync\Network_Layer_IP_Addressing_Revamped.pptx")

COLORS = {
    "bg": RGBColor(11, 32, 59),
    "bg_alt": RGBColor(16, 46, 82),
    "panel": RGBColor(244, 247, 251),
    "title": RGBColor(245, 248, 252),
    "text_light": RGBColor(229, 236, 244),
    "text_dark": RGBColor(38, 50, 66),
    "muted": RGBColor(162, 179, 198),
    "accent": RGBColor(0, 166, 153),
    "accent_2": RGBColor(255, 179, 71),
    "accent_3": RGBColor(80, 142, 255),
}

TEAM = [
    ("Sakshi Negi", "590017526"),
    ("Sahil Singh", "590012382"),
    ("Raghav Patidar", "590013562"),
    ("Harshit Gurja", "590012270"),
    ("Saurav Sahlot", "590014311"),
]


def add_box(slide, left, top, width, height, fill, line_fill=None, radius=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_fill or fill
    return shape


def clear_slide(slide):
    for shape in list(slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)


def set_textbox_text(
    textbox,
    lines,
    color,
    font_size,
    bold=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    space_after=0,
):
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(6)
    frame.margin_bottom = Pt(6)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for index, line in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = line
        para.alignment = align
        para.space_after = Pt(space_after)
        font = para.runs[0].font
        font.name = font_name
        font.size = Pt(font_size if not isinstance(font_size, (list, tuple)) else font_size[index])
        font.bold = bold if not isinstance(bold, (list, tuple)) else bold[index]
        font.color.rgb = color


def add_title_bar(slide, title, subtitle):
    banner = add_box(slide, Inches(0.35), Inches(0.25), Inches(12.6), Inches(1.05), COLORS["bg"])
    banner.line.color.rgb = COLORS["bg"]
    accent = add_box(slide, Inches(0.35), Inches(0.25), Inches(0.22), Inches(1.05), COLORS["accent"])
    accent.line.color.rgb = COLORS["accent"]

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.34), Inches(7.2), Inches(0.42))
    set_textbox_text(title_box, [title], COLORS["title"], 27, bold=True, font_name="Aptos Display")

    subtitle_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.77), Inches(8.8), Inches(0.24))
    set_textbox_text(subtitle_box, [subtitle], COLORS["muted"], 11.5)

    deco_1 = add_box(slide, Inches(11.5), Inches(0.28), Inches(1.0), Inches(0.22), COLORS["accent_2"])
    deco_2 = add_box(slide, Inches(10.95), Inches(0.62), Inches(1.55), Inches(0.18), COLORS["accent_3"])
    deco_3 = add_box(slide, Inches(10.55), Inches(0.92), Inches(1.95), Inches(0.1), COLORS["accent"])
    for shape in (deco_1, deco_2, deco_3):
        shape.line.color.rgb = shape.fill.fore_color.rgb


def add_footer(slide, label):
    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.02), Inches(12.0), Inches(0.24))
    set_textbox_text(footer, [label], COLORS["muted"], 10, align=PP_ALIGN.RIGHT)


def style_slide(slide, index, title, subtitle):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["panel"]

    add_box(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.18), COLORS["accent"])
    add_box(slide, Inches(12.58), Inches(0), Inches(0.75), Inches(7.5), COLORS["bg_alt"])
    add_title_bar(slide, title, subtitle)
    add_footer(slide, f"Data Communication and Networks  |  Slide {index}")


def shape_has_text(shape):
    return hasattr(shape, "text_frame") and shape.has_text_frame


def restyle_existing_text(slide):
    for shape in slide.shapes:
        if not shape_has_text(shape):
            continue
        text = shape.text.strip()
        if not text:
            continue
        for para_index, para in enumerate(shape.text_frame.paragraphs):
            for run in para.runs:
                font = run.font
                if para_index == 0 and font.size and font.size.pt >= 20:
                    font.name = "Aptos Display"
                    font.size = Pt(max(font.size.pt, 22))
                    font.bold = True
                    font.color.rgb = COLORS["bg"]
                else:
                    font.name = "Aptos"
                    if not font.size:
                        font.size = Pt(16)
                    elif font.size.pt > 18:
                        font.size = Pt(min(font.size.pt, 18))
                    font.color.rgb = COLORS["text_dark"]
        if shape.top > Inches(1.4) and shape.left < Inches(12.2):
            shape.fill.background()
            try:
                shape.line.color.rgb = COLORS["panel"]
            except Exception:
                pass


def move_slide(prs, old_index, new_index):
    slides = prs.slides._sldIdLst
    slide = slides[old_index]
    slides.remove(slide)
    slides.insert(new_index, slide)


def clone_slide(prs, slide_index):
    source = prs.slides[slide_index]
    layout = prs.slide_layouts[6]
    dest = prs.slides.add_slide(layout)
    for shape in source.shapes:
        new_el = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return dest


def make_cover(slide):
    clear_slide(slide)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["bg"]

    add_box(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.2), COLORS["accent"])
    add_box(slide, Inches(0), Inches(6.98), Inches(13.33), Inches(0.52), COLORS["accent_2"])
    add_box(slide, Inches(-0.2), Inches(4.95), Inches(13.7), Inches(0.22), COLORS["accent_3"])

    orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.65), Inches(0.55), Inches(3.0), Inches(3.0))
    orb.fill.solid()
    orb.fill.fore_color.rgb = COLORS["bg_alt"]
    orb.line.color.rgb = COLORS["bg_alt"]

    inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.35), Inches(1.25), Inches(1.6), Inches(1.6))
    inner.fill.solid()
    inner.fill.fore_color.rgb = COLORS["accent"]
    inner.line.color.rgb = COLORS["accent"]

    pill = add_box(slide, Inches(0.7), Inches(0.58), Inches(2.3), Inches(0.5), COLORS["accent"])
    pill.line.color.rgb = COLORS["accent"]
    pill_text = slide.shapes.add_textbox(Inches(0.82), Inches(0.67), Inches(2.0), Inches(0.24))
    set_textbox_text(pill_text, ["UNIT IV PRESENTATION"], COLORS["title"], 15, bold=True, font_name="Aptos Display")

    title = slide.shapes.add_textbox(Inches(0.72), Inches(1.45), Inches(7.2), Inches(1.5))
    set_textbox_text(
        title,
        ["Network Layer Design", "& IP Addressing"],
        COLORS["title"],
        [28, 28],
        bold=[True, True],
        font_name="Aptos Display",
        space_after=2,
    )

    subtitle = slide.shapes.add_textbox(Inches(0.75), Inches(2.95), Inches(7.3), Inches(0.8))
    set_textbox_text(
        subtitle,
        [
            "Network Address Translation, IPv4 & IPv6,",
            "Classful/Classless Addressing, Masking & Special Addresses",
        ],
        COLORS["text_light"],
        [16, 16],
    )

    subject_card = add_box(slide, Inches(0.72), Inches(4.08), Inches(4.55), Inches(0.86), COLORS["bg_alt"])
    subject_card.line.color.rgb = COLORS["accent_3"]
    subject_tb = slide.shapes.add_textbox(Inches(0.94), Inches(4.22), Inches(4.1), Inches(0.5))
    set_textbox_text(
        subject_tb,
        ["Subject: Data Communication and Networks"],
        COLORS["title"],
        18,
        bold=True,
        font_name="Aptos Display",
    )

    team_card = add_box(slide, Inches(6.9), Inches(3.65), Inches(5.45), Inches(2.5), COLORS["panel"])
    team_card.line.color.rgb = COLORS["accent_2"]
    team_header = slide.shapes.add_textbox(Inches(7.15), Inches(3.84), Inches(2.2), Inches(0.3))
    set_textbox_text(team_header, ["Prepared By"], COLORS["bg"], 18, bold=True, font_name="Aptos Display")

    left_x = Inches(7.15)
    right_x = Inches(9.7)
    y_positions = [Inches(4.24), Inches(4.63), Inches(5.02)]
    for idx, (name, sap) in enumerate(TEAM):
        x = left_x if idx < 3 else right_x
        y = y_positions[idx] if idx < 3 else y_positions[idx - 3]
        member_box = slide.shapes.add_textbox(x, y, Inches(2.2), Inches(0.34))
        set_textbox_text(member_box, [name], COLORS["text_dark"], 15.5, bold=True)
        sap_box = slide.shapes.add_textbox(x, y + Inches(0.18), Inches(2.2), Inches(0.26))
        set_textbox_text(sap_box, [f"SAP ID: {sap}"], COLORS["bg_alt"], 11.5)

    course = slide.shapes.add_textbox(Inches(0.76), Inches(6.22), Inches(7.6), Inches(0.28))
    set_textbox_text(course, ["Professional classroom deck for tomorrow's presentation"], COLORS["bg"], 12.5, bold=True)


def update_agenda(slide):
    clear_slide(slide)
    style_slide(slide, 2, "Topics Covered", "A polished overview of the unit flow")

    items = [
        "Network Layer Design Issues",
        "Network Address Translation (Concept + Types)",
        "Internet Protocol: IPv4 Addressing",
        "Internet Protocol: IPv6 Addressing",
        "IP Addressing Techniques",
        "Special Addresses and Address Masking",
    ]
    for idx, item in enumerate(items, start=1):
        top = Inches(1.7 + (idx - 1) * 0.82)
        num_box = add_box(slide, Inches(0.95), top, Inches(0.7), Inches(0.48), COLORS["accent"])
        num_tb = slide.shapes.add_textbox(Inches(1.12), top + Inches(0.09), Inches(0.3), Inches(0.2))
        set_textbox_text(num_tb, [f"{idx:02d}"], COLORS["title"], 15, bold=True, font_name="Aptos Display", align=PP_ALIGN.CENTER)
        text_box = add_box(slide, Inches(1.9), top - Inches(0.02), Inches(9.9), Inches(0.52), RGBColor(255, 255, 255), COLORS["accent_3"])
        text_tb = slide.shapes.add_textbox(Inches(2.15), top + Inches(0.07), Inches(9.4), Inches(0.24))
        set_textbox_text(text_tb, [item], COLORS["text_dark"], 17, bold=True)
        text_box.line.width = Pt(1.3)


def add_notes_panel(slide, heading, points):
    panel = add_box(slide, Inches(9.85), Inches(1.55), Inches(2.95), Inches(4.95), RGBColor(255, 255, 255), COLORS["accent"])
    panel.line.width = Pt(1.5)
    head_tb = slide.shapes.add_textbox(Inches(10.08), Inches(1.72), Inches(2.45), Inches(0.35))
    set_textbox_text(head_tb, [heading], COLORS["bg"], 17, bold=True, font_name="Aptos Display")
    body = slide.shapes.add_textbox(Inches(10.08), Inches(2.16), Inches(2.4), Inches(4.0))
    body.text_frame.clear()
    body.text_frame.word_wrap = True
    for i, point in enumerate(points):
        para = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        para.text = point
        para.level = 0
        para.space_after = Pt(8)
        para.bullet = True
        run = para.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(12.5)
        run.font.color.rgb = COLORS["text_dark"]


def tweak_content(slide, index):
    title_map = {
        3: ("Network Layer Design Issues", "Key responsibilities, forwarding behaviour, and service models"),
        4: ("NAT Types & Translation Table", "A closer look at Static NAT, Dynamic NAT, and PAT"),
        5: ("Network Address Translation (NAT)", "Why NAT matters for scalability, security, and IPv4 conservation"),
        6: ("Internet Protocol: IPv4 Addressing", "32-bit addressing structure, notation, and practical understanding"),
        7: ("Internet Protocol: IPv6 Addressing", "128-bit modern addressing for large-scale and future-ready networks"),
        8: ("Classful IP Addressing", "Traditional fixed-size network classes and where they fit"),
        9: ("Classless Addressing (CIDR)", "Flexible prefixes for efficient allocation and routing"),
        10: ("Network & Host Identification", "How subnet masks isolate network bits from host bits"),
        11: ("Special Addresses: Loopback & Broadcast", "Reserved addresses and how they behave inside networks"),
        12: ("Address Masking", "Subnet and wildcard masks with a worked example"),
        13: ("Private IP Ranges & Reserved Blocks", "RFC 1918 spaces and other important non-routable ranges"),
        14: ("Key Takeaways", "Quick recap for a confident presentation ending"),
    }
    if index not in title_map:
        return

    title, subtitle = title_map[index]
    style_slide(slide, index, title, subtitle)
    restyle_existing_text(slide)

    if index == 3:
        add_notes_panel(
            slide,
            "Presenter Focus",
            [
                "Layer 3 decides path selection and forwarding across networks.",
                "Datagram service is flexible; virtual circuits offer predictability.",
                "Store-and-forward allows error checks before forwarding.",
            ],
        )
    elif index == 4:
        add_notes_panel(
            slide,
            "Exam Angle",
            [
                "Static NAT is fixed 1:1 mapping.",
                "Dynamic NAT borrows an address from a public pool.",
                "PAT uses ports so many hosts can share one public IP.",
            ],
        )
    elif index == 5:
        add_notes_panel(
            slide,
            "Why NAT Helps",
            [
                "Saves public IPv4 addresses.",
                "Hides internal addressing from the public Internet.",
                "Lets private networks grow without changing ISP-facing design.",
            ],
        )
    elif index == 6:
        add_notes_panel(
            slide,
            "IPv4 Snapshot",
            [
                "32 bits are written as four decimal octets.",
                "Each octet ranges from 0 to 255.",
                "Address space is limited, so NAT became common.",
            ],
        )
    elif index == 7:
        add_notes_panel(
            slide,
            "IPv6 Snapshot",
            [
                "Huge address space removes address exhaustion pressure.",
                "Hexadecimal notation keeps long addresses manageable.",
                "Supports auto-configuration and smoother end-to-end connectivity.",
            ],
        )
    elif index == 8:
        add_notes_panel(
            slide,
            "Remember",
            [
                "Classes A, B, and C were based on fixed network sizes.",
                "Rigid boundaries wasted address space.",
                "CIDR replaced classful allocation in modern networks.",
            ],
        )
    elif index == 9:
        add_notes_panel(
            slide,
            "CIDR Benefit",
            [
                "Prefix length can match actual need.",
                "Reduces address wastage.",
                "Helps route summarization and smaller routing tables.",
            ],
        )
    elif index == 10:
        add_notes_panel(
            slide,
            "Simple Rule",
            [
                "Mask bit 1 means network part.",
                "Mask bit 0 means host part.",
                "Bitwise AND gives the network ID.",
            ],
        )
    elif index == 11:
        add_notes_panel(
            slide,
            "Use Cases",
            [
                "Loopback checks the local stack without using the network.",
                "Broadcast reaches all devices on a subnet.",
                "Reserved addresses cannot be assigned like normal host IPs.",
            ],
        )
    elif index == 12:
        add_notes_panel(
            slide,
            "Masking Tips",
            [
                "Subnet masks define network boundaries.",
                "Wildcard masks are mostly seen in ACL and routing configs.",
                "Both are inverses of each other.",
            ],
        )
    elif index == 13:
        add_notes_panel(
            slide,
            "Private Blocks",
            [
                "10.0.0.0/8 for very large private networks.",
                "172.16.0.0/12 for medium private blocks.",
                "192.168.0.0/16 is common in homes and labs.",
            ],
        )
    elif index == 14:
        add_notes_panel(
            slide,
            "Closing Line",
            [
                "Start from Layer 3 purpose.",
                "Move to NAT and IP versions.",
                "End with addressing techniques and reserved addresses.",
            ],
        )


def main():
    prs = Presentation(str(SOURCE))
    move_slide(prs, 12, 3)

    make_cover(prs.slides[0])
    update_agenda(prs.slides[1])

    for idx, slide in enumerate(prs.slides, start=1):
        if idx > 2:
            tweak_content(slide, idx)

    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
