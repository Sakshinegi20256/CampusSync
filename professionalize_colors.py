from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path(r"C:\Users\Saksh\Downloads\Network_Layer_IP_Addressing_Updated.pptx")
OUTPUT = Path(r"C:\Users\Saksh\Desktop\CampusSync\Network_Layer_IP_Addressing_Professional.pptx")

PALETTE = {
    "navy": RGBColor(24, 47, 76),
    "navy_soft": RGBColor(34, 63, 100),
    "slate": RGBColor(73, 88, 105),
    "mist": RGBColor(241, 245, 249),
    "white": RGBColor(255, 255, 255),
    "teal": RGBColor(44, 123, 145),
    "teal_soft": RGBColor(215, 234, 239),
    "gold": RGBColor(193, 154, 107),
    "text": RGBColor(37, 46, 56),
    "muted": RGBColor(116, 129, 145),
}

TITLE_MAP = {
    1: ("Network Layer Design & IP Addressing", "Data Communication and Networks"),
    2: ("Topics Covered", "A structured flow for the presentation"),
    3: ("Network Layer Design Issues", "Core responsibilities and routing service choices"),
    4: ("NAT Types & Translation Table", "Static NAT, Dynamic NAT, and PAT at a glance"),
    5: ("Network Address Translation (NAT)", "Private-to-public mapping for scalable Internet access"),
    6: ("Internet Protocol: IPv4 Addressing", "Understanding the 32-bit addressing model"),
    7: ("Internet Protocol: IPv6 Addressing", "Modern 128-bit addressing for future networks"),
    8: ("Classful IP Addressing", "Traditional fixed-size network classes"),
    9: ("Classless Addressing (CIDR)", "Flexible prefixes for efficient allocation"),
    10: ("Network & Host Identification", "Separating network and host bits using masks"),
    11: ("Special Addresses", "Loopback and broadcast behaviour in IP networks"),
    12: ("Address Masking", "Subnet masks and wildcard masks in practice"),
    13: ("Private IP Ranges & Reserved Blocks", "Important non-routable address spaces"),
    14: ("Key Takeaways", "Quick revision before delivery"),
}


def add_textbox(slide, left, top, width, height, text, color, size, bold=False, align=PP_ALIGN.LEFT, font_name="Aptos"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide, idx):
    title, subtitle = TITLE_MAP.get(idx, ("Network Layer Presentation", "Data Communication and Networks"))
    banner = slide.shapes.add_shape(1, Inches(0.35), Inches(0.22), Inches(12.0), Inches(0.92))
    banner.fill.solid()
    banner.fill.fore_color.rgb = PALETTE["navy"]
    banner.line.color.rgb = PALETTE["navy"]

    accent = slide.shapes.add_shape(1, Inches(0.35), Inches(0.22), Inches(0.18), Inches(0.92))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PALETTE["gold"]
    accent.line.color.rgb = PALETTE["gold"]

    add_textbox(slide, Inches(0.62), Inches(0.29), Inches(7.5), Inches(0.32), title, PALETTE["white"], 24, True, font_name="Aptos Display")
    add_textbox(slide, Inches(0.65), Inches(0.64), Inches(7.5), Inches(0.18), subtitle, RGBColor(221, 229, 237), 10.5)

    side = slide.shapes.add_shape(1, Inches(12.66), Inches(0), Inches(0.67), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = PALETTE["navy_soft"]
    side.line.color.rgb = PALETTE["navy_soft"]

    footer = add_textbox(slide, Inches(0.45), Inches(7.03), Inches(12.0), Inches(0.2), f"Data Communication and Networks  |  Slide {idx}", PALETTE["muted"], 9.5, align=PP_ALIGN.RIGHT)
    footer.fill.background()


def restyle_slide(slide, idx):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PALETTE["mist"]

    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PALETTE["teal"]
    top_bar.line.color.rgb = PALETTE["teal"]

    add_header(slide, idx)

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text_frame is None:
            continue
        text = shape.text.strip()
        if not text:
            continue
        if shape.top < Inches(1.2):
            continue

        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = PALETTE["white"]
            shape.line.color.rgb = PALETTE["teal_soft"]
            shape.line.width = Pt(1)
        except Exception:
            pass

        for p_index, para in enumerate(shape.text_frame.paragraphs):
            for run in para.runs:
                font = run.font
                font.name = "Aptos"
                current = font.size.pt if font.size else 14
                if p_index == 0 and current >= 18:
                    font.name = "Aptos Display"
                    font.size = Pt(max(current, 20))
                    font.bold = True
                    font.color.rgb = PALETTE["navy"]
                else:
                    if not font.size:
                        font.size = Pt(13)
                    elif current > 18:
                        font.size = Pt(17)
                    font.color.rgb = PALETTE["text"]


def main():
    prs = Presentation(str(SOURCE))
    for idx, slide in enumerate(prs.slides, start=1):
        restyle_slide(slide, idx)
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
